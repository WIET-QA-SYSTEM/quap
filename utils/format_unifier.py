from pathlib import Path
from pptx import Presentation
from PyPDF2 import PdfFileReader
from tika import parser
from utils.tika_client import TikaClient


class FormatUnifier:
    def __init__(self, host: str, port: int) -> None:
        self._tika_client = TikaClient(host, port)

    def extract_bytes_from_document(self, path_to_file: str) -> str:
        path = Path(path_to_file)
        if not path.exists():
            raise FileExistsError
        if path.suffix == '.pdf':
            pdf = path.open("rb")
            self._tika_client.extract(pdf)
        elif path.suffix == '.pptx':
            return self.extract_from_pptx(path)

    @staticmethod
    def extract_from_pptx(pptx: Path) -> str:
        f = pptx.open("rb")
        prs = Presentation(f)
        text: str = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for t in paragraph.runs:
                        text += ' /n '
                        text += t.text
        return text

'''
f = FormatUnifier()
lal = f.extract_bytes_from_document("D:\Studia\Inne\data\sample.pdf")
print(lal)
'''